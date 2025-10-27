# ======================================================
# certificate_generator_web.py
# SETUP: 
# pip install streamlit pandas python-pptx reportlab Pillow
# streamlit run certificate_generator_web.py --server.address 10.0.0.4
# ======================================================

import streamlit as st
import pandas as pd
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image
import tempfile, os, subprocess, shutil

st.set_page_config(page_title="Certificate Generator", layout="centered")
st.title("Certificate Generator Tool")
st.markdown("Automatically generate certificates using a PowerPoint template.")

# ======================================================
# 1. UPLOAD ATTENDEE LIST
# ======================================================
uploaded_file = st.file_uploader("Upload attendee list (.xlsx, .xls, .csv, .txt)", 
                                 type=["xlsx", "xls", "csv", "txt"])
if uploaded_file:
    if uploaded_file.name.endswith(('.xlsx', '.xls')):
        df = pd.read_excel(uploaded_file)
        names = df.iloc[:, 0].dropna().tolist()
    else:
        df = pd.read_csv(uploaded_file, header=None)
        names = df.iloc[:, 0].dropna().tolist()
    st.success(f"Loaded {len(names)} names.")
else:
    st.stop()

# ======================================================
# 2. EVENT DATE INPUT
# ======================================================
event_date = st.text_input("Enter the event date (e.g., 'October 22, 2025')", "")
if not event_date:
    st.stop()

# ======================================================
# 3. UPLOAD TEMPLATE
# ======================================================
template_file = st.file_uploader("Upload your PowerPoint certificate template (.pptx)", type=["pptx"])
if not template_file:
    st.stop()

try:
    prs = Presentation(template_file)
except Exception as e:
    st.error(f"Error loading PowerPoint file: {e}")
    st.stop()

slide_width_pt = prs.slide_width / 914400 * 72
slide_height_pt = prs.slide_height / 914400 * 72
st.info(f"Template size detected: {slide_width_pt:.1f} x {slide_height_pt:.1f} points")

# ======================================================
# 4. GENERATE PERSONALIZED FILES
# ======================================================
if st.button("Generate Certificates"):

    temp_dir = tempfile.mkdtemp()

    # Save the original template temporarily
    template_path = os.path.join(temp_dir, "template.pptx")
    with open(template_path, "wb") as f:
        f.write(template_file.getbuffer())

    # Generate personalized PPTX
    personalized_files = []
    for idx, name in enumerate(names, start=1):
        prs_copy = Presentation(template_path)
        for slide in prs_copy.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "[NAME]" in run.text or "[DATE]" in run.text:
                                # Backup formatting
                                font = run.font
                                font_name = font.name
                                font_size = font.size
                                font_bold = font.bold
                                font_italic = font.italic
                                font_color_rgb = None
                                try:
                                    if font.color and font.color.rgb:
                                        font_color_rgb = font.color.rgb
                                except AttributeError:
                                    font_color_rgb = None

                                # Replace placeholders
                                run.text = run.text.replace("[NAME]", name).replace("[DATE]", event_date)

                                # Restore formatting
                                run.font.name = font_name
                                run.font.size = font_size
                                run.font.bold = font_bold
                                run.font.italic = font_italic
                                if font_color_rgb:
                                    run.font.color.rgb = font_color_rgb

        cert_file = os.path.join(temp_dir, f"cert_{idx}.pptx")
        prs_copy.save(cert_file)
        personalized_files.append(cert_file)

    st.success("Certificates personalized!")

    # ======================================================
    # 5. CONVERT TO PNG (LibreOffice)
    # ======================================================
    libreoffice_exe = shutil.which("soffice") or shutil.which("libreoffice")
    if not libreoffice_exe:
        potential_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if os.path.exists(potential_path):
            libreoffice_exe = potential_path

    if not libreoffice_exe or not os.path.exists(libreoffice_exe):
        st.error(
            "LibreOffice (soffice.exe) not found.\n"
            "Please ensure LibreOffice is installed and added to your PATH."
        )
        st.stop()

    st.info(f"Using LibreOffice executable: {libreoffice_exe}")
    st.info(f"Generating...")
    
    png_dir = os.path.join(temp_dir, "pngs")
    os.makedirs(png_dir, exist_ok=True)

    for pptx_path in personalized_files:
        result = subprocess.run([
            libreoffice_exe, "--headless", "--convert-to", "png",
            "--outdir", png_dir, pptx_path
        ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

        if result.returncode != 0:
            st.error(f"Error converting {os.path.basename(pptx_path)}:\n{result.stderr}")
            st.stop()

    # ======================================================
    # 6. MERGE PNGs INTO SINGLE PDF
    # ======================================================
    pdf_name = os.path.join(temp_dir, "All_Certificates.pdf")
    c = canvas.Canvas(pdf_name, pagesize=(slide_width_pt, slide_height_pt))

    for png in sorted(os.listdir(png_dir)):
        if png.endswith(".png"):
            img_path = os.path.join(png_dir, png)
            img = Image.open(img_path).convert("RGB")
            img_w, img_h = img.size
            aspect_img = img_w / img_h
            aspect_page = slide_width_pt / slide_height_pt

            if aspect_img > aspect_page:
                scaled_h = slide_width_pt / aspect_img
                offset_y = (slide_height_pt - scaled_h) / 2
                c.drawImage(ImageReader(img), 0, offset_y, width=slide_width_pt, height=scaled_h)
            else:
                scaled_w = slide_height_pt * aspect_img
                offset_x = (slide_width_pt - scaled_w) / 2
                c.drawImage(ImageReader(img), offset_x, 0, width=scaled_w, height=slide_height_pt)

            c.showPage()

    c.save()

    # ======================================================
    # 7. DOWNLOAD BUTTON
    # ======================================================
    with open(pdf_name, "rb") as f:
        st.download_button("Download All Certificates (PDF)", f, file_name="All_Certificates.pdf")

    st.success("All done!")

    # Optional cleanup
    # shutil.rmtree(temp_dir, ignore_errors=True)
